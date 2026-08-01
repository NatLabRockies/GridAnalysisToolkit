"""
Place holder for a file to generate a PowerPoint presentation from figures.
Called at the end of a report or directly from the command line pointed at a directory of figures.
"""

import os
import argparse
from pptx import Presentation
from pptx.util import Inches
from PIL import Image  # To get original image dimensions
from loguru import logger


def create_ppt_from_pngs(root_path, output_file="output.pptx"):
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    current_section = None
    for dirpath, dirnames, filenames in os.walk(root_path):
        rel_path = os.path.relpath(dirpath, root_path)
        section_name = None if rel_path == "." else rel_path.split(os.sep)[0]
        if section_name and section_name != current_section:
            slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            title.text = section_name
            current_section = section_name
            logger.info(f"Added section: {section_name}")
        png_files = [f for f in filenames if f.lower().endswith(".png")]
        for png_file in png_files:
            slide_layout = prs.slide_layouts[5]
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            if title:
                title.text = os.path.splitext(png_file)[0]
            png_path = os.path.join(dirpath, png_file)
            # Get original image dimensions
            with Image.open(png_path) as img:
                orig_width, orig_height = img.size
            aspect_ratio = orig_width / orig_height
            # Target area: 10 × 6 inches
            width = Inches(10)
            height = width / aspect_ratio

            left = Inches(0)  # Center in 8-inch width
            top = Inches(1.5)  # Keep title space

            max_height = Inches(6.5)

            if height > max_height:
                height = max_height
                width = height * aspect_ratio
                left = Inches(10) - width

            try:
                slide.shapes.add_picture(png_path, left, top, width, height)
                logger.info(f"Added {png_file} to presentation (w={width}, h={height})")
            except Exception as e:
                logger.warning(f"Error adding {png_file}: {str(e)}")
    prs.save(output_file)
    logger.info(f"Presentation saved as {output_file}")


def main():
    parser = argparse.ArgumentParser(
        description="Create a PowerPoint presentation from PNG files in a directory"
    )
    parser.add_argument("root_directory", help="Root directory to search for PNG files")
    parser.add_argument(
        "-o",
        "--output",
        default="output.pptx",
        help="Output PowerPoint filename (default: output.pptx)",
    )
    args = parser.parse_args()
    create_ppt_from_pngs(args.root_directory, args.output)


if __name__ == "__main__":
    main()
